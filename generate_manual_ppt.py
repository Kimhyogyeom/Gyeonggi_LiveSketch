# -*- coding: utf-8 -*-
"""
경기 LiveSketch 운영 매뉴얼 PPT 생성 스크립트
실행: python generate_manual_ppt.py
출력: LiveSketch_운영매뉴얼.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ────────────────────────────────────────────────
#  색상 팔레트
# ────────────────────────────────────────────────
C_BG_DARK   = RGBColor(0x0D, 0x1B, 0x2A)   # 진남색 배경
C_BG_CARD   = RGBColor(0x1A, 0x2E, 0x44)   # 카드 배경
C_ACCENT    = RGBColor(0x00, 0xBF, 0xFF)   # 청록 강조
C_GREEN     = RGBColor(0x2E, 0xCC, 0x71)   # 초록 (정상/성공)
C_ORANGE    = RGBColor(0xFF, 0x8C, 0x00)   # 주황 (주의)
C_RED       = RGBColor(0xFF, 0x3B, 0x3B)   # 빨강 (오류)
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)   # 흰색
C_GRAY      = RGBColor(0xAA, 0xBB, 0xCC)   # 연한 회색
C_YELLOW    = RGBColor(0xFF, 0xD7, 0x00)   # 노랑 (강조)

# ────────────────────────────────────────────────
#  슬라이드 크기 (16:9 와이드)
# ────────────────────────────────────────────────
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


# ────────────────────────────────────────────────
#  헬퍼 함수
# ────────────────────────────────────────────────

def add_filled_rect(slide, left, top, width, height, fill_color, alpha=None):
    """단색 사각형 추가"""
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_text_box(slide, text, left, top, width, height,
                 font_size=18, bold=False, color=C_WHITE,
                 align=PP_ALIGN.LEFT, wrap=True, italic=False):
    """텍스트 박스 추가"""
    from pptx.util import Pt
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "맑은 고딕"
    return txBox


def add_bg(slide, color=C_BG_DARK):
    """슬라이드 배경 채우기"""
    add_filled_rect(slide, 0, 0, SLIDE_W, SLIDE_H, color)


def add_header_bar(slide, title, subtitle=None,
                   bar_color=C_ACCENT, title_color=C_BG_DARK):
    """상단 헤더 바 추가"""
    bar_h = Inches(1.1)
    add_filled_rect(slide, 0, 0, SLIDE_W, bar_h, bar_color)
    add_text_box(slide, title,
                 Inches(0.4), Inches(0.12), Inches(10), Inches(0.7),
                 font_size=30, bold=True, color=title_color)
    if subtitle:
        add_text_box(slide, subtitle,
                     Inches(0.4), Inches(0.78), Inches(10), Inches(0.35),
                     font_size=13, bold=False, color=RGBColor(0x22,0x44,0x66))


def step_box(slide, num_text, desc_lines, left, top, width=Inches(5.8),
             num_color=C_ACCENT, bg_color=C_BG_CARD, note=None):
    """번호 + 설명 카드 박스"""
    box_h = Inches(1.0 + 0.35 * max(0, len(desc_lines) - 1))
    add_filled_rect(slide, left, top, width, box_h, bg_color)

    # 번호 뱃지
    badge_w = Inches(0.6)
    add_filled_rect(slide, left, top, badge_w, box_h, num_color)
    add_text_box(slide, num_text,
                 left, top + Inches(0.2),
                 badge_w, Inches(0.5),
                 font_size=22, bold=True,
                 color=C_BG_DARK, align=PP_ALIGN.CENTER)

    # 설명 텍스트
    desc = "\n".join(desc_lines)
    add_text_box(slide, desc,
                 left + badge_w + Inches(0.1),
                 top + Inches(0.1),
                 width - badge_w - Inches(0.15),
                 box_h - Inches(0.05),
                 font_size=15, bold=False, color=C_WHITE)
    if note:
        add_text_box(slide, f"※ {note}",
                     left + badge_w + Inches(0.1),
                     top + box_h - Inches(0.3),
                     width - badge_w - Inches(0.2),
                     Inches(0.3),
                     font_size=11, italic=True, color=C_GRAY)
    return box_h


def error_card(slide, icon, title, causes, solutions,
               left, top, width=Inches(5.9)):
    """오류 카드"""
    card_h = Inches(0.45 + 0.28 * (len(causes) + len(solutions) + 1.5))
    add_filled_rect(slide, left, top, width, card_h, C_BG_CARD)
    # 좌측 강조선
    add_filled_rect(slide, left, top, Inches(0.06), card_h, C_RED)

    # 제목
    add_text_box(slide, f"{icon}  {title}",
                 left + Inches(0.15), top + Inches(0.08),
                 width - Inches(0.2), Inches(0.35),
                 font_size=15, bold=True, color=C_RED)

    y = top + Inches(0.42)
    for c in causes:
        add_text_box(slide, f"▸ {c}",
                     left + Inches(0.2), y,
                     width - Inches(0.25), Inches(0.28),
                     font_size=12, color=C_ORANGE)
        y += Inches(0.27)

    y += Inches(0.05)
    for s in solutions:
        add_text_box(slide, f"✔ {s}",
                     left + Inches(0.2), y,
                     width - Inches(0.25), Inches(0.28),
                     font_size=12, color=C_GREEN)
        y += Inches(0.27)
    return card_h


# ════════════════════════════════════════════════
#  슬라이드 생성
# ════════════════════════════════════════════════

prs = new_prs()
blank_layout = prs.slide_layouts[6]  # blank


# ──────────────────────────────────────────────
#  슬라이드 1 ─ 표지
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, C_BG_DARK)

# 그라디언트 느낌을 위한 상단 선
add_filled_rect(slide, 0, 0, SLIDE_W, Inches(0.08), C_ACCENT)

# 중앙 제목 블록
add_filled_rect(slide, Inches(1.5), Inches(2.0), Inches(10.3), Inches(3.3), C_BG_CARD)
add_filled_rect(slide, Inches(1.5), Inches(2.0), Inches(0.12), Inches(3.3), C_ACCENT)

add_text_box(slide, "경기 LiveSketch",
             Inches(1.9), Inches(2.15), Inches(9.5), Inches(1.0),
             font_size=44, bold=True, color=C_ACCENT, align=PP_ALIGN.LEFT)
add_text_box(slide, "전시 운영 매뉴얼",
             Inches(1.9), Inches(3.05), Inches(9.5), Inches(0.7),
             font_size=28, bold=False, color=C_WHITE, align=PP_ALIGN.LEFT)
add_text_box(slide, "시작 방법  ·  종료 방법  ·  오류 대처",
             Inches(1.9), Inches(3.7), Inches(9.5), Inches(0.5),
             font_size=16, bold=False, color=C_GRAY, align=PP_ALIGN.LEFT)

add_text_box(slide, "경기도 희귀 생물 체험 전시",
             Inches(0), Inches(6.8), SLIDE_W, Inches(0.4),
             font_size=12, color=C_GRAY, align=PP_ALIGN.CENTER)

add_filled_rect(slide, 0, Inches(7.42), SLIDE_W, Inches(0.08), C_ACCENT)


# ──────────────────────────────────────────────
#  슬라이드 2 ─ 목차
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_header_bar(slide, "목차", "이 매뉴얼에서 다루는 내용")

sections = [
    ("01", "시스템 개요",    "LiveSketch가 어떻게 작동하는지 간단히 설명합니다.",    C_ACCENT),
    ("02", "시작 방법",      "전시 운영 전 준비 사항과 단계별 시작 절차입니다.",    C_GREEN),
    ("03", "정상 작동 확인", "프로그램이 올바르게 동작하는지 확인하는 방법입니다.", C_ACCENT),
    ("04", "종료 방법",      "전시 종료 후 안전하게 시스템을 종료하는 방법입니다.", C_ORANGE),
    ("05", "오류 대처",      "자주 발생하는 문제와 해결 방법을 안내합니다.",        C_RED),
]

for i, (num, title, desc, col) in enumerate(sections):
    y = Inches(1.35) + i * Inches(1.1)
    add_filled_rect(slide, Inches(0.4), y, Inches(12.5), Inches(0.95), C_BG_CARD)
    add_filled_rect(slide, Inches(0.4), y, Inches(0.7), Inches(0.95), col)
    add_text_box(slide, num,
                 Inches(0.4), y + Inches(0.2), Inches(0.7), Inches(0.5),
                 font_size=18, bold=True, color=C_BG_DARK, align=PP_ALIGN.CENTER)
    add_text_box(slide, title,
                 Inches(1.25), y + Inches(0.08), Inches(3.5), Inches(0.45),
                 font_size=18, bold=True, color=col)
    add_text_box(slide, desc,
                 Inches(1.25), y + Inches(0.52), Inches(11.0), Inches(0.35),
                 font_size=13, color=C_GRAY)


# ──────────────────────────────────────────────
#  슬라이드 3 ─ 시스템 개요
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_header_bar(slide, "01  시스템 개요", "LiveSketch 작동 흐름")

# 흐름 화살표
steps_flow = [
    ("①", "학습지 스캔",    "스캐너로\n학습지 스캔"),
    ("②", "QR 인식",        "QR 코드로\n캐릭터 식별"),
    ("③", "색상 추출",      "채색 부분\n색상 분석"),
    ("④", "캐릭터 등장",    "화면에\n3D 캐릭터 생성"),
    ("⑤", "배너 표시",      "캐릭터 이름\n배너 팝업"),
]

box_w = Inches(2.2)
total_w = len(steps_flow) * box_w + (len(steps_flow) - 1) * Inches(0.15)
start_x = (SLIDE_W - total_w) / 2

for i, (num, title, desc) in enumerate(steps_flow):
    x = start_x + i * (box_w + Inches(0.15))
    y = Inches(1.7)
    add_filled_rect(slide, x, y, box_w, Inches(2.2), C_BG_CARD)
    add_filled_rect(slide, x, y, box_w, Inches(0.45), C_ACCENT)
    add_text_box(slide, num,
                 x, y + Inches(0.05), box_w, Inches(0.38),
                 font_size=16, bold=True, color=C_BG_DARK, align=PP_ALIGN.CENTER)
    add_text_box(slide, title,
                 x, y + Inches(0.52), box_w, Inches(0.4),
                 font_size=14, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    add_text_box(slide, desc,
                 x, y + Inches(0.95), box_w, Inches(1.1),
                 font_size=13, color=C_WHITE, align=PP_ALIGN.CENTER)
    if i < len(steps_flow) - 1:
        ax = x + box_w + Inches(0.01)
        add_text_box(slide, "▶",
                     ax, y + Inches(0.8), Inches(0.14), Inches(0.5),
                     font_size=18, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

# 하단 설명
add_text_box(slide,
             "학습지에 그려진 그림의 색상이 자동으로 화면 속 캐릭터에 입혀집니다.\n"
             "스캐너 → PC → 화면까지 모두 자동으로 처리되므로, 운영자는 스캐너 작동만 확인하면 됩니다.",
             Inches(0.5), Inches(4.2), Inches(12.3), Inches(1.1),
             font_size=14, color=C_GRAY)

# 스캔 폴더 정보 박스
add_filled_rect(slide, Inches(0.5), Inches(5.4), Inches(12.3), Inches(1.6), C_BG_CARD)
add_filled_rect(slide, Inches(0.5), Inches(5.4), Inches(0.06), Inches(1.6), C_YELLOW)
add_text_box(slide, "📁  주요 폴더 위치",
             Inches(0.65), Inches(5.5), Inches(8), Inches(0.35),
             font_size=14, bold=True, color=C_YELLOW)
add_text_box(slide,
             "스캔 이미지 저장 폴더:   C:\\ProgramData\\LiveSketch\\Scans\n"
             "배경 이미지 폴더:          프로그램 옆  BackgroundImage\\Use\\",
             Inches(0.75), Inches(5.85), Inches(11.8), Inches(0.9),
             font_size=13, color=C_WHITE)


# ──────────────────────────────────────────────
#  슬라이드 4 ─ 시작 방법 (준비 체크리스트)
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_header_bar(slide, "02  시작 방법", "운영 전 준비 사항 체크리스트", bar_color=C_GREEN,
               title_color=C_BG_DARK)

checklist = [
    ("✅", "스캐너 전원",   "스캐너 전원이 켜져 있고 PC와 USB 연결되어 있는지 확인"),
    ("✅", "PC 전원",       "모니터(전시 화면)와 PC 전원이 켜져 있는지 확인"),
    ("✅", "스캔 소프트웨어", "스캐너 전용 소프트웨어가 열려 있는지 확인 (EPSON Scan 등)"),
    ("✅", "배경 이미지",   "BackgroundImage\\Use\\ 폴더에 배경 이미지가 있는지 확인"),
    ("✅", "학습지 준비",   "QR 코드가 인쇄된 채색 학습지가 준비되어 있는지 확인"),
]

for i, (icon, name, desc) in enumerate(checklist):
    y = Inches(1.35) + i * Inches(1.0)
    add_filled_rect(slide, Inches(0.4), y, Inches(12.5), Inches(0.85), C_BG_CARD)
    add_text_box(slide, icon,
                 Inches(0.5), y + Inches(0.15), Inches(0.5), Inches(0.5),
                 font_size=22, align=PP_ALIGN.CENTER, color=C_GREEN)
    add_text_box(slide, name,
                 Inches(1.15), y + Inches(0.08), Inches(2.5), Inches(0.38),
                 font_size=16, bold=True, color=C_GREEN)
    add_text_box(slide, desc,
                 Inches(1.15), y + Inches(0.46), Inches(11.5), Inches(0.33),
                 font_size=13, color=C_GRAY)

add_text_box(slide, "위 항목을 모두 확인했으면 다음 슬라이드의 시작 절차를 따라주세요.",
             Inches(0.4), Inches(6.5), Inches(12.5), Inches(0.5),
             font_size=13, italic=True, color=C_GRAY, align=PP_ALIGN.CENTER)


# ──────────────────────────────────────────────
#  슬라이드 5 ─ 시작 방법 (단계별 실행)
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_header_bar(slide, "02  시작 방법", "단계별 실행 절차", bar_color=C_GREEN,
               title_color=C_BG_DARK)

left_steps = [
    ("1", ["LiveSketch 프로그램 실행",
           "바탕화면의 'LiveSketch' 아이콘을 더블클릭합니다."]),
    ("2", ["배경 화면 확인",
           "전시 화면에 배경 이미지가 슬라이드쇼로 나타나면 정상입니다."]),
    ("3", ["스캐너 연결 확인",
           "스캐너 소프트웨어에서 스캐너가 '준비됨' 상태인지 확인합니다."]),
]
right_steps = [
    ("4", ["학습지 스캔",
           "QR 코드가 오른쪽 상단에 오도록 학습지를 스캐너 유리판에 올립니다.",
           "스캔 버튼을 눌러 스캔합니다."]),
    ("5", ["캐릭터 등장 확인",
           "전시 화면에 캐릭터 배너가 뜨고 잠시 후 캐릭터가 나타납니다.",
           "색상이 올바르게 입혀지면 성공입니다!"]),
]

y = Inches(1.35)
for step in left_steps:
    h = step_box(slide, step[0], step[1], Inches(0.4), y, width=Inches(6.0), bg_color=C_BG_CARD)
    y += h + Inches(0.12)

y = Inches(1.35)
for step in right_steps:
    h = step_box(slide, step[0], step[1], Inches(6.9), y, width=Inches(6.0), bg_color=C_BG_CARD)
    y += h + Inches(0.12)

add_text_box(slide,
             "🔑  QR 코드가 오른쪽 상단이 되도록 학습지를 놓는 것이 가장 중요합니다!",
             Inches(0.4), Inches(6.5), Inches(12.5), Inches(0.6),
             font_size=14, bold=True, color=C_YELLOW, align=PP_ALIGN.CENTER)


# ──────────────────────────────────────────────
#  슬라이드 6 ─ 정상 작동 확인
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_header_bar(slide, "03  정상 작동 확인", "이렇게 보이면 올바르게 작동하는 것입니다")

checks = [
    ("🖼️", "배경 슬라이드쇼",
     "전시 화면에 배경 이미지가 일정 간격으로 자연스럽게 전환됩니다."),
    ("📢", "캐릭터 배너",
     "스캔 후 오른쪽에서 캐릭터 이름 배너가 부드럽게 나타났다 사라집니다."),
    ("🐸", "캐릭터 등장",
     "배너 후 3D 캐릭터가 화면에 나타나고 자연스럽게 움직입니다."),
    ("🎨", "색상 적용",
     "학습지에 칠한 색상이 화면 속 캐릭터에 그대로 입혀집니다."),
    ("🔊", "효과음",
     "캐릭터 등장 시 효과음이 재생됩니다. (스피커 볼륨 확인)"),
]

for i, (icon, title, desc) in enumerate(checks):
    col = i % 2
    row = i // 2
    x = Inches(0.4) + col * Inches(6.5)
    y = Inches(1.35) + row * Inches(1.65)
    w = Inches(6.0)
    add_filled_rect(slide, x, y, w, Inches(1.45), C_BG_CARD)
    add_filled_rect(slide, x, y, Inches(0.06), Inches(1.45), C_GREEN)
    add_text_box(slide, icon + "  " + title,
                 x + Inches(0.2), y + Inches(0.1), w, Inches(0.4),
                 font_size=16, bold=True, color=C_GREEN)
    add_text_box(slide, desc,
                 x + Inches(0.2), y + Inches(0.55), w - Inches(0.3), Inches(0.75),
                 font_size=13, color=C_WHITE)

# 마지막 항목 (5번째, 가운데)
i = 4
x = Inches(0.4) + 2 * Inches(3.25) + Inches(0.0)  # 중앙
x = (SLIDE_W - Inches(6.0)) / 2
y = Inches(1.35) + 2 * Inches(1.65)
w = Inches(6.0)
add_filled_rect(slide, x, y, w, Inches(1.45), C_BG_CARD)
add_filled_rect(slide, x, y, Inches(0.06), Inches(1.45), C_GREEN)
icon, title, desc = checks[4]
add_text_box(slide, icon + "  " + title,
             x + Inches(0.2), y + Inches(0.1), w, Inches(0.4),
             font_size=16, bold=True, color=C_GREEN)
add_text_box(slide, desc,
             x + Inches(0.2), y + Inches(0.55), w - Inches(0.3), Inches(0.75),
             font_size=13, color=C_WHITE)


# ──────────────────────────────────────────────
#  슬라이드 7 ─ 종료 방법
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_header_bar(slide, "04  종료 방법", "전시 종료 후 안전한 시스템 종료 순서",
               bar_color=C_ORANGE, title_color=C_BG_DARK)

end_steps = [
    ("1", ["LiveSketch 프로그램 종료",
           "화면 오른쪽 상단의 ✕ 버튼 클릭 또는 Alt + F4 키를 누릅니다."],
     None),
    ("2", ["스캐너 소프트웨어 종료",
           "스캐너 전용 소프트웨어를 닫습니다."],
     None),
    ("3", ["스캐너 전원 끄기",
           "스캐너 전원 버튼을 눌러 끕니다."],
     "USB 케이블은 뽑지 않아도 됩니다."),
    ("4", ["PC 종료",
           "시작 → 전원 → 시스템 종료를 선택합니다."],
     "화면이 꺼진 것을 확인한 후 자리를 비우세요."),
]

y = Inches(1.35)
for step in end_steps:
    note = step[2] if len(step) > 2 else None
    h = step_box(slide, step[0], step[1], Inches(1.5), y,
                 width=Inches(10.3), num_color=C_ORANGE, bg_color=C_BG_CARD, note=note)
    y += h + Inches(0.18)

add_filled_rect(slide, Inches(0.4), Inches(6.25), Inches(12.5), Inches(0.85), C_BG_CARD)
add_filled_rect(slide, Inches(0.4), Inches(6.25), Inches(0.06), Inches(0.85), C_YELLOW)
add_text_box(slide,
             "⚠️  스캐너 유리판에 학습지가 남아있지 않은지 확인 후 종료하세요.",
             Inches(0.6), Inches(6.38), Inches(12.0), Inches(0.6),
             font_size=14, bold=True, color=C_YELLOW)


# ──────────────────────────────────────────────
#  슬라이드 8 ─ 오류 대처 (QR / 캐릭터 미등장)
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_header_bar(slide, "05  오류 대처 ①", "QR 코드 인식 실패  /  캐릭터가 나타나지 않음",
               bar_color=C_RED, title_color=C_WHITE)

# 왼쪽 카드
error_card(slide,
           "🔲", "QR 코드를 인식하지 못함 (캐릭터가 안 나옴)",
           causes=[
               "학습지 방향이 잘못됨 (QR이 오른쪽 상단이 아님)",
               "QR 코드 위에 낙서/오염이 있음",
               "스캔 해상도가 너무 낮음",
           ],
           solutions=[
               "학습지를 QR 코드가 오른쪽 상단이 되도록 다시 놓기",
               "QR 코드 주변을 깨끗하게 닦기",
               "스캐너 소프트웨어에서 해상도 300 DPI 이상으로 설정",
               "다시 스캔 시도",
           ],
           left=Inches(0.4), top=Inches(1.35), width=Inches(6.0))

# 오른쪽 카드
error_card(slide,
           "🐾", "캐릭터가 화면에 등장하지 않음",
           causes=[
               "프로그램이 스캔 파일을 아직 처리 중",
               "스캔 폴더 경로가 잘못 설정됨",
               "스캐너 소프트웨어가 지정 폴더에 저장 안 함",
           ],
           solutions=[
               "스캔 후 5~10초 대기",
               "스캐너 저장 폴더 확인: C:\\ProgramData\\LiveSketch\\Scans",
               "폴더에 최근 JPG 파일이 생성되는지 확인",
               "프로그램 재시작 후 다시 스캔",
           ],
           left=Inches(6.9), top=Inches(1.35), width=Inches(6.0))

# 공통 TIP
add_filled_rect(slide, Inches(0.4), Inches(5.6), Inches(12.5), Inches(1.5), C_BG_CARD)
add_filled_rect(slide, Inches(0.4), Inches(5.6), Inches(0.06), Inches(1.5), C_ACCENT)
add_text_box(slide, "💡  공통 해결 방법",
             Inches(0.6), Inches(5.65), Inches(8), Inches(0.38),
             font_size=14, bold=True, color=C_ACCENT)
add_text_box(slide,
             "1. 프로그램을 완전히 종료(Alt+F4) 후 다시 실행    "
             "2. PC를 재시작하면 대부분의 문제가 해결됩니다.",
             Inches(0.7), Inches(6.05), Inches(12.0), Inches(0.85),
             font_size=13, color=C_WHITE)


# ──────────────────────────────────────────────
#  슬라이드 9 ─ 오류 대처 (색상 이상 / 배경 이상)
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_header_bar(slide, "05  오류 대처 ②", "색상이 이상함  /  배경 이미지가 바뀌지 않음",
               bar_color=C_RED, title_color=C_WHITE)

error_card(slide,
           "🎨", "캐릭터 색상이 이상하게 입혀짐",
           causes=[
               "유리판에 먼지/지문이 있음",
               "학습지가 기울어져 스캔됨",
               "색이 너무 연하게 칠해짐",
           ],
           solutions=[
               "스캐너 유리판을 부드러운 천으로 닦기",
               "학습지를 유리판 모서리에 맞춰 바르게 놓기",
               "색을 좀 더 진하게 채색한 학습지로 다시 시도",
               "다시 스캔",
           ],
           left=Inches(0.4), top=Inches(1.35), width=Inches(6.0))

error_card(slide,
           "🖼️", "배경 이미지가 변경되지 않음 / 검은 화면",
           causes=[
               "BackgroundImage\\Use\\ 폴더가 비어 있음",
               "지원하지 않는 이미지 형식 사용",
               "파일 이름에 한글/특수문자 포함",
           ],
           solutions=[
               "BackgroundImage\\Use\\ 폴더에 JPG 또는 PNG 이미지 추가",
               "파일 이름을 영문/숫자로 변경 (예: bg01.jpg)",
               "프로그램 재시작",
           ],
           left=Inches(6.9), top=Inches(1.35), width=Inches(6.0))

add_filled_rect(slide, Inches(0.4), Inches(5.0), Inches(12.5), Inches(2.1), C_BG_CARD)
add_filled_rect(slide, Inches(0.4), Inches(5.0), Inches(0.06), Inches(2.1), C_YELLOW)
add_text_box(slide, "📁  배경 이미지 관리 방법",
             Inches(0.6), Inches(5.08), Inches(8), Inches(0.38),
             font_size=14, bold=True, color=C_YELLOW)
add_text_box(slide,
             "• 보관 폴더:  BackgroundImage\\Resource\\  (여기에 원본 저장)\n"
             "• 표시 폴더:  BackgroundImage\\Use\\  (여기 있는 이미지만 화면에 표시)\n"
             "• Resource에서 Use로 이미지를 복사/이동하면 다음 실행부터 표시됩니다.",
             Inches(0.7), Inches(5.5), Inches(12.0), Inches(1.45),
             font_size=13, color=C_WHITE)


# ──────────────────────────────────────────────
#  슬라이드 10 ─ 오류 대처 (프로그램 응답없음 / 재시작)
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_header_bar(slide, "05  오류 대처 ③", "프로그램 응답 없음  /  긴급 재시작",
               bar_color=C_RED, title_color=C_WHITE)

# 왼쪽: 강제종료 방법
add_filled_rect(slide, Inches(0.4), Inches(1.35), Inches(5.8), Inches(4.7), C_BG_CARD)
add_filled_rect(slide, Inches(0.4), Inches(1.35), Inches(0.06), Inches(4.7), C_RED)
add_text_box(slide, "🔴  프로그램 강제 종료 방법",
             Inches(0.6), Inches(1.45), Inches(5.5), Inches(0.38),
             font_size=15, bold=True, color=C_RED)

force_quit_steps = [
    "Ctrl + Alt + Del  키를 동시에 누릅니다.",
    "'작업 관리자 열기'를 클릭합니다.",
    "목록에서 'LiveSketch' 를 찾아 클릭합니다.",
    "오른쪽 클릭 → '작업 끝내기' 를 선택합니다.",
    "바탕화면으로 돌아온 뒤 프로그램을 다시 실행합니다.",
]
for i, step in enumerate(force_quit_steps):
    y_step = Inches(1.95) + i * Inches(0.6)
    add_filled_rect(slide, Inches(0.6), y_step, Inches(0.4), Inches(0.42),
                    C_RED)
    add_text_box(slide, str(i+1),
                 Inches(0.6), y_step + Inches(0.03), Inches(0.4), Inches(0.38),
                 font_size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, step,
                 Inches(1.15), y_step + Inches(0.05),
                 Inches(4.85), Inches(0.42),
                 font_size=13, color=C_WHITE)

# 오른쪽: PC 완전 재시작
add_filled_rect(slide, Inches(6.7), Inches(1.35), Inches(6.2), Inches(4.7), C_BG_CARD)
add_filled_rect(slide, Inches(6.7), Inches(1.35), Inches(0.06), Inches(4.7), C_ORANGE)
add_text_box(slide, "🔄  PC 완전 재시작 (최후 수단)",
             Inches(6.9), Inches(1.45), Inches(5.9), Inches(0.38),
             font_size=15, bold=True, color=C_ORANGE)

restart_info = [
    ("강제종료가 안 될 때", "PC 전원 버튼을 5초간 길게 눌러 강제 종료"),
    ("재시작 방법", "시작 → 전원 → 다시 시작 선택"),
    ("부팅 후", "LiveSketch 프로그램을 다시 실행"),
    ("여전히 문제 시", "담당자에게 연락 (다음 페이지 참조)"),
]
for i, (title, desc) in enumerate(restart_info):
    y_item = Inches(1.95) + i * Inches(0.9)
    add_text_box(slide, title,
                 Inches(6.9), y_item, Inches(5.8), Inches(0.35),
                 font_size=13, bold=True, color=C_ORANGE)
    add_text_box(slide, desc,
                 Inches(6.9), y_item + Inches(0.35), Inches(5.8), Inches(0.42),
                 font_size=13, color=C_GRAY)

# 하단 경고
add_filled_rect(slide, Inches(0.4), Inches(6.25), Inches(12.5), Inches(0.85), C_BG_CARD)
add_filled_rect(slide, Inches(0.4), Inches(6.25), Inches(0.06), Inches(0.85), C_RED)
add_text_box(slide,
             "⚠️  PC를 강제 종료한 경우, 다음 실행 시 디스크 검사가 진행될 수 있습니다. "
             "검사가 끝날 때까지 기다려 주세요.",
             Inches(0.6), Inches(6.35), Inches(12.1), Inches(0.65),
             font_size=13, color=C_YELLOW)


# ──────────────────────────────────────────────
#  슬라이드 11 ─ 오류 요약 표
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_header_bar(slide, "05  오류 대처 ─ 요약표", "빠른 참조용")

headers = ["증상", "원인", "해결 방법"]
rows = [
    ["캐릭터가 안 나옴", "QR 방향 오류", "학습지 방향 확인 (QR = 오른쪽 상단)"],
    ["캐릭터가 안 나옴", "스캔 폴더 문제", "저장 폴더 확인 후 재스캔"],
    ["색상이 이상함",    "유리판 오염",   "유리판 닦고 다시 스캔"],
    ["색상이 이상함",    "기울어진 스캔", "학습지 바르게 놓고 재스캔"],
    ["배경이 없음",      "이미지 폴더 비어 있음", "Use 폴더에 JPG 이미지 추가"],
    ["프로그램 멈춤",    "메모리/오류",   "작업 관리자로 강제 종료 후 재실행"],
]

col_widths = [Inches(3.2), Inches(3.5), Inches(5.5)]
col_x = [Inches(0.4), Inches(3.7), Inches(7.3)]
row_h = Inches(0.65)
header_y = Inches(1.35)

# 헤더 행
for j, (hd, cw, cx) in enumerate(zip(headers, col_widths, col_x)):
    add_filled_rect(slide, cx, header_y, cw, row_h, C_ACCENT)
    add_text_box(slide, hd,
                 cx + Inches(0.1), header_y + Inches(0.1),
                 cw - Inches(0.15), row_h - Inches(0.15),
                 font_size=15, bold=True, color=C_BG_DARK, align=PP_ALIGN.CENTER)

# 데이터 행
row_colors = [C_BG_CARD, RGBColor(0x14, 0x24, 0x38)]
for i, row in enumerate(rows):
    ry = header_y + row_h + i * row_h
    bg = row_colors[i % 2]
    for j, (cell, cw, cx) in enumerate(zip(row, col_widths, col_x)):
        add_filled_rect(slide, cx, ry, cw, row_h, bg)
        col_c = C_ORANGE if j == 1 else (C_GREEN if j == 2 else C_WHITE)
        add_text_box(slide, cell,
                     cx + Inches(0.1), ry + Inches(0.1),
                     cw - Inches(0.15), row_h - Inches(0.1),
                     font_size=13, color=col_c)

add_text_box(slide,
             "위 방법으로 해결되지 않으면 담당 기술자에게 연락하세요.",
             Inches(0.4), Inches(6.7), Inches(12.5), Inches(0.45),
             font_size=13, italic=True, color=C_GRAY, align=PP_ALIGN.CENTER)


# ──────────────────────────────────────────────
#  슬라이드 12 ─ 마무리 / 담당자 연락처
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_filled_rect(slide, 0, 0, SLIDE_W, Inches(0.08), C_ACCENT)

add_text_box(slide, "이 매뉴얼을 참고하여 원활한 전시 운영 부탁드립니다.",
             Inches(0), Inches(1.6), SLIDE_W, Inches(0.7),
             font_size=20, color=C_GRAY, align=PP_ALIGN.CENTER)
add_text_box(slide, "경기 LiveSketch",
             Inches(0), Inches(2.3), SLIDE_W, Inches(1.0),
             font_size=42, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
add_text_box(slide, "전시 운영 매뉴얼",
             Inches(0), Inches(3.25), SLIDE_W, Inches(0.6),
             font_size=24, color=C_WHITE, align=PP_ALIGN.CENTER)

# 연락처 박스 (빈칸으로 남겨 운영자가 직접 작성)
add_filled_rect(slide, Inches(3.2), Inches(4.2), Inches(6.9), Inches(2.3), C_BG_CARD)
add_filled_rect(slide, Inches(3.2), Inches(4.2), Inches(0.06), Inches(2.3), C_ACCENT)
add_text_box(slide, "📞  담당 기술자 연락처",
             Inches(3.4), Inches(4.3), Inches(6.5), Inches(0.38),
             font_size=14, bold=True, color=C_ACCENT)
add_text_box(slide,
             "이름:  ___________________________\n"
             "전화:  ___________________________\n"
             "이메일: __________________________",
             Inches(3.5), Inches(4.75), Inches(6.3), Inches(1.55),
             font_size=14, color=C_GRAY)

add_filled_rect(slide, 0, Inches(7.42), SLIDE_W, Inches(0.08), C_ACCENT)


# ────────────────────────────────────────────────
#  저장
# ────────────────────────────────────────────────
output_path = "LiveSketch_운영매뉴얼.pptx"
prs.save(output_path)
import sys, os
sys.stdout.reconfigure(encoding='utf-8') if hasattr(sys.stdout, 'reconfigure') else None
print(f"[완료] 저장: {output_path}  |  슬라이드 수: {len(prs.slides)}장")
